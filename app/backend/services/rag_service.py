import asyncio
import hashlib
from pathlib import Path

from docx import Document as DocxDocument
from langchain_core.documents import Document as LangchainDocument
from langchain_openai import OpenAIEmbeddings
from langchain_qdrant import QdrantVectorStore
from langchain_text_splitters import RecursiveCharacterTextSplitter
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader
from qdrant_client import QdrantClient, models
from qdrant_client.models import Distance, VectorParams
from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession

from app.backend.config import get_settings
from app.backend.database.models import IndexedFile
from app.backend.database.session import AsyncSessionLocal
from app.backend.logger import get_logger

logger = get_logger(__name__)


class RAGService:
    def __init__(self):
        self.settings = get_settings()
        self.collection_name = self.settings.COLLECTION_NAME
        self.qdrant_url = self.settings.QDRANT_URL
        self.vector_store: QdrantVectorStore | None = None

    def search_in_file(self, file_path: str, question: str) -> str:
        """Search within a specific file - auto-indexes if needed"""
        try:
            path = Path(file_path).expanduser().resolve()

            if not path.exists():
                return f"Error: File '{file_path}' not found"

            if not path.is_file():
                return f"Error: '{file_path}' is not a file"

            async def check_and_index():
                async with AsyncSessionLocal() as db:
                    result = await db.execute(
                        select(IndexedFile).where(IndexedFile.file_path == str(path))
                    )
                    existing = result.scalar_one_or_none()
                    file_hash = self._calculate_hash(str(path))

                    if not existing or existing.file_hash != file_hash:
                        logger.info(f"Auto-indexing {path}...")
                        await self._index_document(str(path), db)

            asyncio.run(check_and_index())

            self._init_vector_store()
            results = self.vector_store.similarity_search( # type: ignore
                question,
                k=3,
                filter=models.Filter(
                    must=[
                        models.FieldCondition(
                            key="metadata.source",
                            match=models.MatchValue(value=str(path)),
                        )
                    ]
                ),
            )

            if not results:
                return f"No relevant information found in {file_path}"

            response = f"Found {len(results)} relevant sections in {file_path}:\n\n"
            for i, doc in enumerate(results, 1):
                response += f"Section {i}:\n{doc.page_content}\n\n"

            return response.strip()

        except Exception as e:
            logger.error(f"Failed to search in {file_path}: {e}")
            return f"Error: {str(e)}"

    def _init_vector_store(self):
        """Initialize vector store connection"""
        if self.vector_store:
            return
        client = QdrantClient(url=self.qdrant_url)
        self.vector_store = QdrantVectorStore(
            client=client,
            collection_name=self.collection_name,
            embedding=OpenAIEmbeddings(),
        )

    async def _index_document(self, document_path: str, db: AsyncSession) -> None:
        """Index document with database tracking and recursive chunking"""
        path = Path(document_path)
        file_hash = self._calculate_hash(document_path)

        result = await db.execute(select(IndexedFile).where(IndexedFile.file_path == str(path)))
        existing = result.scalar_one_or_none()

        if existing and existing.file_hash == file_hash:
            logger.info(f"File unchanged, skipping: {path}")
            return

        text_content = self._load_file(document_path)
        doc = LangchainDocument(page_content=text_content, metadata={"source": str(path)})

        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            separators=["\n\n", "\n", ". ", " ", ""],
        )
        chunks = text_splitter.split_documents([doc])

        embeddings = OpenAIEmbeddings()
        client = QdrantClient(url=self.settings.QDRANT_URL)

        try:
            client.create_collection(
                collection_name=self.settings.COLLECTION_NAME,
                vectors_config=VectorParams(size=1536, distance=Distance.COSINE),
            )
            logger.info(f"Created collection '{self.settings.COLLECTION_NAME}'")
        except Exception as e:
            logger.debug(f"Collection '{self.settings.COLLECTION_NAME}' already exists: {e}")

        vector_store = QdrantVectorStore(
            client=client,
            collection_name=self.settings.COLLECTION_NAME,
            embedding=embeddings,
        )
        vector_store.add_documents(chunks)

        # Track in database
        if existing:
            existing.file_hash = file_hash
            existing.file_size = path.stat().st_size
        else:
            db.add(
                IndexedFile(
                    file_path=str(path),
                    file_hash=file_hash,
                    file_type=path.suffix,
                    file_size=path.stat().st_size,
                )
            )

        await db.commit()
        logger.info(f"Indexed {len(chunks)} chunks from '{path}'")

    def _calculate_hash(self, file_path: str) -> str:
        """Calculate MD5 hash of file to detect changes"""
        with open(file_path, "rb") as f:
            return hashlib.md5(f.read()).hexdigest()

    def _load_file(self, file_path: str) -> str:
        """Load text from various file types"""
        path = Path(file_path)
        suffix = path.suffix.lower()

        if suffix == ".docx":
            doc = DocxDocument(file_path)
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

        elif suffix == ".pdf":
            reader = PdfReader(file_path)
            text = []
            for page in reader.pages:
                if page_text := page.extract_text():
                    text.append(page_text)
            return "\n\n".join(text) if text else ""

        elif suffix == ".xlsx":
            workbook = load_workbook(file_path, data_only=True)
            text = []
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text.append(f"Sheet: {sheet_name}\n")
                for row in sheet.iter_rows(values_only=True):
                    row_text = "\t".join(str(cell) if cell is not None else "" for cell in row)
                    if row_text.strip():
                        text.append(row_text)
            return "\n".join(text) if text else ""

        elif suffix == ".pptx":
            prs = Presentation(file_path)
            text = []
            for i, slide in enumerate(prs.slides, 1):
                text.append(f"Slide {i}:")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text.append(shape.text)
            return "\n\n".join(text) if text else ""

        else:
            try:
                return path.read_text(encoding="utf-8")
            except UnicodeDecodeError as e:
                raise ValueError(
                    f"File is not a supported format or has unsupported encoding: {file_path}"
                ) from e
